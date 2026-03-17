"""
Add a Service Architecture slide to aaa-platform-presentation.pptx
Slide 9144000 x 5143500 EMU (10" x 5.625")
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = r"C:\Users\rony.idan\source\public-repo\aaa-cloud-native\docs\aaa-platform-presentation.pptx"

# Colors
NAVY       = RGBColor(0x1F, 0x38, 0x64)
BLUE       = RGBColor(0x44, 0x72, 0xC4)
BLUE_DARK  = RGBColor(0x2E, 0x5E, 0xB8)
BLUE_LIGHT = RGBColor(0x9D, 0xB6, 0xE4)
GREEN      = RGBColor(0x70, 0xAD, 0x47)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_COLOR   = RGBColor(0xEE, 0xF3, 0xFB)

W = 9144000   # slide width
H = 5143500   # slide height

prs = Presentation(PPTX_PATH)

# --- Add blank slide using the blank layout ---
blank_layout = None
for layout in prs.slide_layouts:
    if layout.name in ("Blank", "blank"):
        blank_layout = layout
        break
if blank_layout is None:
    blank_layout = prs.slide_layouts[6]  # fallback to index 6 (usually blank)

slide = prs.slides.add_slide(blank_layout)
shapes = slide.shapes

# Remove any default placeholders from blank layout
for ph in list(slide.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

# --- Helper: solid fill on shape ---
def solid_fill(shape, rgb):
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:pattFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")

def no_line(shape):
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for tag in ("a:solidFill", "a:gradFill", "a:noFill"):
        el = ln.find(qn(tag))
        if el is not None:
            ln.remove(el)
    etree.SubElement(ln, qn("a:noFill"))

def set_line(shape, rgb, width_pt=0.75):
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(width_pt * 12700)))
    for tag in ("a:solidFill", "a:gradFill", "a:noFill"):
        el = ln.find(qn(tag))
        if el is not None:
            ln.remove(el)
    solidFill = etree.SubElement(ln, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")

def add_label(shape, text, font_size_pt=9, bold=False, color=WHITE, align=PP_ALIGN.CENTER, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    sp = shape._element
    txBody = sp.find(qn("p:txBody"))
    if txBody is not None:
        bodyPr = txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("anchor", "ctr")

# ============================================================
# BACKGROUND
# ============================================================
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = BG_COLOR

# ============================================================
# HEADER BAR
# ============================================================
header = shapes.add_shape(1, Emu(0), Emu(0), Emu(W), Emu(480000))
solid_fill(header, NAVY)
no_line(header)
tf = header.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Service Architecture"
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = WHITE
txBody = header._element.find(qn("p:txBody"))
bodyPr = txBody.find(qn("a:bodyPr"))
bodyPr.set("anchor", "ctr")
bodyPr.set("lIns", str(Emu(320000)))

# ============================================================
# LAYOUT CONSTANTS  (all in EMU)
# ============================================================
Y_DB  = 860000
Y_SVC = 1700000
Y_RAD = 2780000
Y_SMF = 3860000

H_DB  = 300000
W_DB  = 980000

H_OVL = 360000
W_OVL_LG = 1650000

H_RAD = 300000
W_RAD = 1380000

H_SMF = 250000
W_SMF = 570000

H_UI  = 360000
W_UI  = 360000

# ============================================================
# HELPER: cylinder shape
# ============================================================
def add_cylinder(cx, cy, w, h, label):
    cap_h = int(h * 0.38)
    body_h = h - cap_h // 2
    x = cx - w // 2
    y = cy - h // 2

    # Body rectangle
    body = shapes.add_shape(1, Emu(x), Emu(y + cap_h//2), Emu(w), Emu(body_h))
    solid_fill(body, BLUE)
    no_line(body)

    # Bottom cap
    bot = shapes.add_shape(9, Emu(x), Emu(y + body_h - cap_h//3), Emu(w), Emu(cap_h))
    solid_fill(bot, BLUE_DARK)
    set_line(bot, BLUE, 0.5)

    # Top cap
    top = shapes.add_shape(9, Emu(x), Emu(y), Emu(w), Emu(cap_h))
    solid_fill(top, BLUE_LIGHT)
    set_line(top, BLUE, 0.75)

    add_label(top, label, font_size_pt=8, bold=True, color=WHITE)
    return (cx, cy - h//2)  # top connection point

# ============================================================
# HELPER: oval
# ============================================================
def add_oval(cx, cy, w, h, label, font_size=9):
    x = cx - w // 2
    y = cy - h // 2
    shape = shapes.add_shape(9, Emu(x), Emu(y), Emu(w), Emu(h))
    solid_fill(shape, BLUE)
    set_line(shape, BLUE_LIGHT, 0.75)
    add_label(shape, label, font_size_pt=font_size, bold=True, color=WHITE)
    return shape

# ============================================================
# HELPER: rounded rect
# ============================================================
def add_rrect(cx, cy, w, h, label, bg_color=BLUE, font_size=7.5):
    x = cx - w // 2
    y = cy - h // 2
    shape = shapes.add_shape(5, Emu(x), Emu(y), Emu(w), Emu(h))
    solid_fill(shape, bg_color)
    set_line(shape, bg_color, 0.5)
    add_label(shape, label, font_size_pt=font_size, bold=True, color=WHITE)
    return shape

# ============================================================
# HELPER: arrow connector
# ============================================================
_arrow_id = [100]
def add_arrow(x1, y1, x2, y2):
    _arrow_id[0] += 1
    spTree = slide.shapes._spTree
    # Determine geometry: normal or flipped
    flip = ""
    if x1 > x2 and y1 <= y2:
        flip = 'flipH="1"'
    elif x1 <= x2 and y1 > y2:
        flip = 'flipV="1"'
    elif x1 > x2 and y1 > y2:
        flip = 'flipH="1" flipV="1"'

    ox = min(x1, x2)
    oy = min(y1, y2)
    cx = abs(x2 - x1)
    cy_val = abs(y2 - y1)
    if cx == 0:
        cx = 1
    if cy_val == 0:
        cy_val = 1

    cxn_xml = f"""<p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvCxnSpPr>
    <p:cNvPr id="{_arrow_id[0]}" name="arrow{_arrow_id[0]}"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm {flip}>
      <a:off x="{ox}" y="{oy}"/>
      <a:ext cx="{cx}" cy="{cy_val}"/>
    </a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="9525">
      <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
      <a:headEnd type="none"/>
      <a:tailEnd type="arrow" w="sm" len="sm"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>"""
    cxn_el = etree.fromstring(cxn_xml)
    spTree.append(cxn_el)

# ============================================================
# ROW 1: PostgreSQL cylinders
# ============================================================
db_xs = [2200000, 4050000, 5900000]
for cx in db_xs:
    add_cylinder(cx, Y_DB, W_DB, H_DB, "PostgreSQL")

# ============================================================
# ROW 2: Services
# ============================================================
CX_LOOKUP = 2100000
# Shadow oval
add_oval(CX_LOOKUP + 130000, Y_SVC + 70000, W_OVL_LG, H_OVL, "", font_size=1)
# Main oval
add_oval(CX_LOOKUP, Y_SVC, W_OVL_LG, H_OVL, "Lookup service", font_size=9)

CX_CONN   = 5100000
W_CONN    = 2500000
add_oval(CX_CONN, Y_SVC, W_CONN, H_OVL, "1st Connection &\nProvisioning", font_size=8)

# UI circle
CX_UI = 8200000
ui_shape = shapes.add_shape(9, Emu(CX_UI - W_UI//2), Emu(Y_SVC - H_UI//2), Emu(W_UI), Emu(H_UI))
solid_fill(ui_shape, BLUE)
set_line(ui_shape, WHITE, 0.75)
add_label(ui_shape, "UI", font_size_pt=10, bold=True, color=WHITE)

# ============================================================
# ROW 3: Radius to Rest
# ============================================================
rad_xs = [1500000, 3900000, 6300000]
for cx in rad_xs:
    add_oval(cx, Y_RAD, W_RAD, H_RAD, "Radius to Rest", font_size=8)

# ============================================================
# ROW 4: SMF + QA
# ============================================================
smf_left   = [440000, 1040000, 1640000, 2240000]
smf_center = [3100000, 3700000, 4300000, 4900000]
smf_right  = [5750000, 6350000]

for cx in smf_left + smf_center + smf_right:
    add_rrect(cx, Y_SMF, W_SMF, H_SMF, "SMF")

add_rrect(7500000, Y_SMF, 860000, H_SMF, "QA Automation", bg_color=GREEN, font_size=8)

# ============================================================
# ARROWS: DB → Services
# ============================================================
# Each DB connects to both Lookup and 1st Connection
for db_cx in db_xs:
    add_arrow(db_cx, Y_DB + H_DB//2, CX_LOOKUP, Y_SVC - H_OVL//2)
    add_arrow(db_cx, Y_DB + H_DB//2, CX_CONN,   Y_SVC - H_OVL//2)

# ============================================================
# ARROWS: Services → Radius
# ============================================================
# Lookup → all 3 Radius
for rad_cx in rad_xs:
    add_arrow(CX_LOOKUP, Y_SVC + H_OVL//2, rad_cx, Y_RAD - H_RAD//2)

# 1st Connection → all 3 Radius
for rad_cx in rad_xs:
    add_arrow(CX_CONN, Y_SVC + H_OVL//2, rad_cx, Y_RAD - H_RAD//2)

# ============================================================
# ARROWS: Radius → SMF
# ============================================================
for cx in smf_left:
    add_arrow(rad_xs[0], Y_RAD + H_RAD//2, cx, Y_SMF - H_SMF//2)

for cx in smf_center:
    add_arrow(rad_xs[1], Y_RAD + H_RAD//2, cx, Y_SMF - H_SMF//2)

for cx in smf_right:
    add_arrow(rad_xs[2], Y_RAD + H_RAD//2, cx, Y_SMF - H_SMF//2)

# ============================================================
# ARROW: UI → 1st Connection
# ============================================================
add_arrow(CX_UI - W_UI//2, Y_SVC, CX_CONN + W_CONN//2, Y_SVC)

# ============================================================
# SAVE
# ============================================================
prs.save(PPTX_PATH)
print(f"Saved to: {PPTX_PATH}")
print(f"Total slides: {len(prs.slides)}")
